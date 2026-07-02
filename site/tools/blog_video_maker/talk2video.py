#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
talk2video.py — 把「PPTX/PDF + 每页中英讲稿」离线合成为带 Kokoro 配音 + 中英双语字幕的 MP4

设计要点
  • 中文居顶(纯黑+极细浅蓝描边)、英文居底(黑字+极细粉描边)；各留固定安全边距，
    文字过长向正文方向溢出，仅描边不铺底 → 溢出处画面透明可见(BorderStyle=1)。
  • 字幕特效 --effect: karaoke (卡拉OK) / handwrite (逐字书写) / particle(粒子消散) / none；
    中英共用同一时间轴 → 双语同步、且与语音同步。
  • 画质：幻灯片按 DPI 渲染到「不低于输出分辨率」，再用 lanczos 下采样(超采样抗锯齿)；
    内存只随单张图线性增长(非按帧/时长) → 提 DPI 几乎不增内存。
  • 响度：每段统一 loudnorm 到 -12 LUFS(比流媒体 -16 约响 4dB)。

依赖: ffmpeg(含 libass)、poppler-utils(pdftoppm)、PPT 还需 libreoffice；
      Python: kokoro torch numpy (PPTX 视频叠加还需 python-pptx)
字体: fonts/SourceSans3-Regular.otf  fonts/SourceHanSansSC-Regular.otf

python assets/media/blog/script/talk2video.py assets/media/blog/script \
       --input assets/media/blog/script/GE-Sim.pptx --effect handwrite
"""

import argparse, json, os, re, shutil, subprocess, sys, threading, wave
from concurrent.futures import ThreadPoolExecutor
from dataclasses import dataclass
from pathlib import Path
from typing import List, Tuple
import numpy as np

os.environ["HF_HUB_OFFLINE"] = "1"
os.environ["TRANSFORMERS_OFFLINE"] = "1"
os.environ.setdefault("HF_HUB_DISABLE_TELEMETRY", "1")

SAMPLE_RATE = 24000
KOKORO_REPO = "hexgrad/Kokoro-82M"
FONT_EXT = {".otf", ".ttf", ".ttc", ".otc"}
VEXT = {"mp4", "mov", "m4v", "webm", "avi", "wmv", "mpg", "mpeg"}
ANIM = {"anim", "animClr", "animEffect", "animMotion", "animRot", "animScale"}
R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"

VOICES_HELP = (
    "Kokoro 音色(默认 am_michael)：\n"
    "  美式(a) 男 am_adam/am_michael/am_onyx  女 af_heart/af_bella/af_nicole\n"
    "  英式(b) 男 bm_george/bm_lewis        女 bf_emma/bf_isabella\n"
    "  其他需配合 --lang：中文(z) zf_xiaoxiao/zm_yunyang  日(j) jf_alpha  法(f) ff_siwis")


# ============================== 配置 ==============================
@dataclass
class Config:
    project: Path
    voice: str = "am_michael"
    lang_code: str = "a"
    speed: float = 1.0
    jobs: int = 8
    preset: str = "medium"
    crf: int = 18
    device: str = "auto"
    effect: str = "karaoke"             # karaoke | handwrite | particle | none
    loudness: float = -12.0             # 统一响度(LUFS)，越大越响(-12≈比-16响4dB)
    # —— 停顿(秒) ——
    lead_silence: float = 0.20
    tail_silence: float = 0.60
    sentence_pause: float = 0.25
    min_slide: float = 1.5
    # —— 视频 ——
    width: int = 1920
    height: int = 1080
    fps: int = 30
    dpi: int = 200                      # 见文末「DPI×内存」分析
    # —— 字幕 ——
    fonts_dir: Path = None
    eng_file: Path = None
    han_file: Path = None
    margin_h: int = 60                  # 左右安全边距
    shadow: float = 0.0                 # 字幕阴影(px)，0=无
    # 英文(底部)：黑字 + 极细粉描边
    eng_font: str = "Source Sans 3"
    eng_size: int = 46
    eng_margin: int = 40               # 距画面底部固定安全高度
    eng_outline: float = 2.0           # 粉描边尽可能细(1px=清晰下限)
    eng_primary: str = "&H00000000"       # 黑(卡拉OK已读)
    eng_second: str = "&H008C8C8C"        # 灰(卡拉OK未读)
    eng_border: str = "&H00C1B6FF"        # 浅粉描边
    # 中文(顶部)：纯黑 + 极细浅蓝描边
    han_font: str = "Source Han Sans SC"
    han_size: int = 4
    han_margin: int = 4
    han_outline: float = 3.0
    han_primary: str = "&H00000000"       # 纯黑
    han_second: str = "&H008C8C8C"
    han_border: str = "&H00E6D8AD"        # 浅蓝描边

    def __post_init__(self):
        self.project = Path(self.project).resolve()
        self.fonts_dir = (Path(self.fonts_dir).resolve() if self.fonts_dir
                          else Path(__file__).resolve().parent / "fonts")

    @property
    def build(self) -> Path:
        return self.project / "build"


@dataclass
class Vid:
    path: Path
    x: int; y: int; w: int; h: int


@dataclass
class Cue:
    """一句字幕：en/zh 文本 + 语音起止 [t0, t1](秒)。"""
    en: str
    zh: str
    t0: float
    t1: float


# ============================== 工具 ==============================
def run(cmd):
    cmd = [str(c) for c in cmd]
    r = subprocess.run(cmd, stdout=subprocess.DEVNULL, stderr=subprocess.PIPE, text=True)
    if r.returncode:
        raise subprocess.CalledProcessError(r.returncode, cmd, stderr=r.stderr)
    return r


def need(tool, hint):
    if not shutil.which(tool):
        sys.exit(f"[ERROR] 缺少 '{tool}'，请安装：{hint}")


def ensure_libass():
    try:
        out = subprocess.run(["ffmpeg", "-hide_banner", "-filters"],
                             capture_output=True, text=True).stdout
    except Exception:
        return
    if not re.search(r"\bass\b", out):
        sys.exit("[ERROR] 当前 ffmpeg 未编译 libass(缺 ass 滤镜)，无法烧录字幕。")


class Bar:
    """线程安全的轻量进度条。"""
    def __init__(self, total, desc="处理"):
        self.t = max(1, total); self.n = 0; self.desc = desc
        self.lock = threading.Lock(); self._draw()

    def update(self, k=1):
        with self.lock:
            self.n = min(self.t, self.n + k); self._draw()

    def _draw(self):
        w = 30; f = int(w * self.n / self.t)
        sys.stdout.write(f"\r{self.desc} |{'█'*f}{'░'*(w-f)}| {self.n}/{self.t}")
        sys.stdout.flush()
        if self.n >= self.t:
            sys.stdout.write("\n"); sys.stdout.flush()


def natural_key(p: Path):
    return [int(n) for n in re.findall(r"\d+", p.stem)] or [0]


def _find_font(d: Path, *keys):
    if not d.exists():
        return None
    for p in sorted(d.glob("*")):
        if p.suffix.lower() in FONT_EXT:
            n = re.sub(r"[^a-z0-9]", "", p.stem.lower())
            if all(k in n for k in keys):
                return p
    return None


def resolve_font(d: Path, family: str, *fallbacks):
    fam = re.sub(r"[^a-z0-9]", "", family.lower())
    return _find_font(d, fam) or next((f for k in fallbacks if (f := _find_font(d, k))), None)


# ====================== 输入定位 / 渲染 PNG ======================
def find_deck(project: Path, explicit=None) -> Path:
    if explicit:
        p = Path(explicit)
        if not p.is_absolute() and not p.exists():
            p = project / p
        p = p.resolve()
        if not p.exists():
            sys.exit(f"[ERROR] 输入文件不存在：{explicit}")
        if p.suffix.lower() not in (".pptx", ".ppt", ".pdf"):
            sys.exit(f"[ERROR] 不支持的输入类型 {p.suffix}")
        return p
    decks = (sorted(project.glob("*.pptx")) + sorted(project.glob("*.ppt"))
             + sorted(project.glob("*.pdf")))
    if not decks:
        sys.exit(f"[ERROR] {project} 下找不到 .pptx/.ppt/.pdf(可用 --input 指定)")
    if len(decks) > 1:
        print(f"[warn] 多个候选输入，使用：{decks[0].name}")
    return decks[0]


def deck_to_images(cfg: Config, deck: Path) -> List[Path]:
    slides = cfg.build / "slides"; slides.mkdir(parents=True, exist_ok=True)
    for f in slides.glob("slide-*.png"):
        f.unlink()
    need("pdftoppm", "sudo apt-get install -y poppler-utils")
    if deck.suffix.lower() == ".pdf":
        pdf = deck
    else:
        soffice = shutil.which("libreoffice") or shutil.which("soffice")
        if not soffice:
            sys.exit("[ERROR] 缺少 LibreOffice：sudo apt-get install -y libreoffice")
        run([soffice, "--headless", "--norestore",
             "-env:UserInstallation=file:///tmp/talk2video_lo",
             "--convert-to", "pdf", "--outdir", cfg.build, deck])
        pdf = cfg.build / (deck.stem + ".pdf")
        if not pdf.exists():
            sys.exit("[ERROR] 幻灯片 -> PDF 失败")
    # 关键：按 DPI 渲染到「不低于输出」，后续 ffmpeg 用 lanczos 下采样(超采样抗锯齿)
    run(["pdftoppm", "-png", "-r", cfg.dpi, pdf, slides / "slide"])
    imgs = sorted(slides.glob("slide-*.png"), key=natural_key)
    if not imgs:
        sys.exit("[ERROR] 没有生成任何幻灯片图片")
    print(f"[slides] 渲染完成，共 {len(imgs)} 页 @ {cfg.dpi}DPI")
    return imgs


# ====================== 讲稿读取(JSON / TXT) ======================
_ABBR = re.compile(r"\b(e\.g|i\.e|etc|vs|Mr|Mrs|Ms|Dr|Inc|Ltd|No)\.", re.I)

def split_sentences(text: str) -> List[str]:
    text = re.sub(r"\s+", " ", text.strip())
    if not text:
        return []
    text = _ABBR.sub(lambda m: m.group(0).replace(".", "<DOT>"), text)
    return [p.replace("<DOT>", ".").strip()
            for p in re.split(r"(?<=[.!?])\s+", text) if p.strip()]


def load_script(path: Path) -> List[Tuple[str, str]]:
    if path.suffix.lower() == ".json":
        data = json.loads(path.read_text(encoding="utf-8"))
        sents = data["sentences"] if isinstance(data, dict) else data
        out = []
        for s in sents:
            if isinstance(s, dict):
                en, zh = s.get("en", ""), s.get("zh", "")
            elif isinstance(s, (list, tuple)):
                en, zh = s[0], (s[1] if len(s) > 1 else "")
            else:
                en, zh = str(s), ""
            en, zh = en.strip(), (zh or "").strip()
            if en:
                out.append((en, zh))
        return out
    return [(s, "") for s in split_sentences(path.read_text(encoding="utf-8"))]


def find_scripts(project: Path) -> List[Path]:
    talk = project / "talk"
    return (sorted(talk.glob("*.json"), key=natural_key)
            or sorted(talk.glob("*.txt"), key=natural_key))


# ============================ TTS 合成 ============================
def to_numpy(a) -> np.ndarray:
    if hasattr(a, "detach"):
        a = a.detach().cpu().numpy()
    return np.asarray(a, dtype=np.float32)


def silence(sec: float) -> np.ndarray:
    return np.zeros(max(0, round(sec * SAMPLE_RATE)), dtype=np.float32)


def make_pipeline(cfg: Config):
    import torch
    from kokoro import KPipeline
    device = "cuda" if (cfg.device == "auto" and torch.cuda.is_available()) \
        else ("cpu" if cfg.device == "auto" else cfg.device)
    try:
        pipe = KPipeline(lang_code=cfg.lang_code, repo_id=KOKORO_REPO, device=device)
    except TypeError:
        pipe = KPipeline(lang_code=cfg.lang_code, repo_id=KOKORO_REPO)
    print(f"[tts] device={device} · voice={cfg.voice} · speed={cfg.speed} · 离线")
    return pipe


def synth(pipe, text: str, cfg: Config) -> np.ndarray:
    chunks = [to_numpy(a) for _, _, a in pipe(text, voice=cfg.voice, speed=cfg.speed)]
    chunks = [c for c in chunks if c.size]
    return np.concatenate(chunks) if chunks else silence(0.0)


def synthesize_slide(pipe, sentences, cfg: Config):
    """逐句合成并记录语音起止 [t0,t1]；句间留 sentence_pause。返回 (audio, cues, total)。"""
    pieces, cues, n = [], [], 0

    def push(s):
        nonlocal n
        pieces.append(s); n += len(s)

    def now():
        return n / SAMPLE_RATE

    push(silence(cfg.lead_silence))
    for i, (en, zh) in enumerate(sentences):
        t0 = now()
        a = synth(pipe, en, cfg)
        if a.size == 0:
            a = silence(0.10)
        push(a)
        t1 = now()
        cues.append(Cue(en, zh, t0, t1))
        if i < len(sentences) - 1:
            push(silence(cfg.sentence_pause))
    push(silence(cfg.tail_silence))

    audio = np.concatenate(pieces) if pieces else silence(0.0)
    pad = cfg.min_slide - len(audio) / SAMPLE_RATE
    if pad > 0:
        audio = np.concatenate([audio, silence(pad)])
    peak = float(np.max(np.abs(audio))) if audio.size else 0.0
    if peak > 0:                                          # 峰值归一 → 拉满量程(再交 loudnorm 定响度)
        audio = audio * (0.97 / peak)
    return audio, cues, len(audio) / SAMPLE_RATE


def write_wav(path: Path, audio: np.ndarray):
    pcm = (np.clip(audio, -1.0, 1.0) * 32767.0).astype("<i2")
    with wave.open(str(path), "wb") as w:
        w.setnchannels(1); w.setsampwidth(2); w.setframerate(SAMPLE_RATE)
        w.writeframes(pcm.tobytes())


# ============================ ASS 字幕 ============================
def _ass_time(t: float) -> str:
    cs = int(round(max(0.0, t) * 100))
    h, cs = divmod(cs, 360000); m, cs = divmod(cs, 6000); s, cs = divmod(cs, 100)
    return f"{h:d}:{m:02d}:{s:02d}.{cs:02d}"


def _ass_text(s: str) -> str:
    s = (s or "").strip().replace("\\", "")
    s = s.replace("{", "(").replace("}", ")")
    return re.sub(r"\s+", " ", s)


def _units(text: str, cjk: bool, fine: bool) -> List[str]:
    """切分动画单元：中文按字；英文 fine=逐字符(含空格)，否则按词(含尾随空格)。"""
    if cjk:
        return list(text) or [text]
    if fine:
        return list(text) or [text]
    ws = text.split(" ")
    return [w + (" " if i < len(ws) - 1 else "") for i, w in enumerate(ws) if w] or [text]


def _cum_cs(units, total_cs):
    """按字符数把 total_cs(厘秒) 分配给各单元，返回每个单元结束累计时刻。"""
    wts = [max(1, len(u.strip())) for u in units]
    sw = sum(wts) or 1
    out, acc = [], 0
    for w in wts:
        acc += w
        out.append(round(total_cs * acc / sw))
    return out


def _fx_body(text: str, t0: float, t1: float, end: float, effect: str, cjk: bool) -> str:
    """生成一行字幕正文(含特效内联标签)。语音段[t0,t1]、整体显示到 end；中英共用时间窗→同步。"""
    if effect == "none" or not text:
        return text
    dur_cs = max(1, round((t1 - t0) * 100))

    if effect == "karaoke":                       # 词/字级填充扫光：未读灰 → 已读黑
        units = _units(text, cjk, fine=False)
        prev, parts = 0, []
        for u, e in zip(units, _cum_cs(units, dur_cs)):
            parts.append(f"{{\\kf{max(0, e - prev)}}}{u}"); prev = e
        return "".join(parts)

    if effect == "handwrite":                     # 逐字符按语速淡入(近似书写)
        units = _units(text, cjk, fine=True)
        prev, parts = 0, []
        for u, e in zip(units, _cum_cs(units, dur_cs)):
            parts.append(f"{{\\alpha&HFF&\\t({prev*10},{e*10},\\alpha&H00&)}}{u}"); prev = e
        return "".join(parts)

    # particle：语音段全程可见，[t1,end] 内逐单元 模糊+放大+淡出 → 消散
    units = _units(text, cjk, fine=False)
    base = (t1 - t0)
    diss = min(0.8, max(0.05, end - t1))
    n = len(units); parts = []
    for i, u in enumerate(units):
        s = round((base + diss * i / n) * 1000)
        e = round((base + diss * (i + 1) / n) * 1000)
        parts.append(f"{{\\t({s},{e},\\alpha&HFF&\\blur8\\fscx150\\fscy150)}}{u}")
    return "".join(parts)


def build_ass(cues: List[Cue], cfg: Config, total: float) -> str:
    W, H, m = cfg.width, cfg.height, cfg.margin_h
    out = [
        "[Script Info]", "ScriptType: v4.00+", f"PlayResX: {W}", f"PlayResY: {H}",
        "WrapStyle: 0", "ScaledBorderAndShadow: yes", "YCbCr Matrix: TV.709", "",
        "[V4+ Styles]",
        ("Format: Name, Fontname, Fontsize, PrimaryColour, SecondaryColour, OutlineColour, "
         "BackColour, Bold, Italic, Underline, StrikeOut, ScaleX, ScaleY, Spacing, Angle, "
         "BorderStyle, Outline, Shadow, Alignment, MarginL, MarginR, MarginV, Encoding"),
        # 英文居底(Alignment 2)：单行贴底，溢出自动向上堆叠进正文；仅描边不铺底→透明
        (f"Style: En,{cfg.eng_font},{cfg.eng_size},{cfg.eng_primary},{cfg.eng_second},"
         f"{cfg.eng_border},&H00000000,0,0,0,0,100,100,0,0,1,{cfg.eng_outline},{cfg.shadow},"
         f"2,{m},{m},{cfg.eng_margin},1"),
        # 中文居顶(Alignment 8)：单行贴顶，溢出自动向下堆叠进正文；仅描边不铺底→透明
        (f"Style: Zh,{cfg.han_font},{cfg.han_size},{cfg.han_primary},{cfg.han_second},"
         f"{cfg.han_border},&H00000000,0,0,0,0,100,100,0,0,1,{cfg.han_outline},{cfg.shadow},"
         f"8,{m},{m},{cfg.han_margin},1"),
        "", "[Events]",
        "Format: Layer, Start, End, Style, Name, MarginL, MarginR, MarginV, Effect, Text",
    ]
    for i, c in enumerate(cues):
        en, zh = _ass_text(c.en), _ass_text(c.zh)
        end = cues[i + 1].t0 if i + 1 < len(cues) else total
        if en:
            out.append(f"Dialogue: 0,{_ass_time(c.t0)},{_ass_time(end)},En,,0,0,0,,"
                       + _fx_body(en, c.t0, c.t1, end, cfg.effect, False))
        if zh:
            out.append(f"Dialogue: 0,{_ass_time(c.t0)},{_ass_time(end)},Zh,,0,0,0,,"
                       + _fx_body(zh, c.t0, c.t1, end, cfg.effect, True))
    return "\n".join(out) + "\n"


# ====================== PPTX 内嵌视频提取 ======================
def extract_pptx_videos(cfg: Config, deck: Path):
    """返回 (videos_per{0基页->[Vid]}, 含动画页数, 幻灯片总数)。坐标已换算到 W×H 画面。"""
    from pptx import Presentation
    prs = Presentation(str(deck))
    sw, sh, W, H = prs.slide_width, prs.slide_height, cfg.width, cfg.height
    if sw / sh >= W / H:
        dw, dh = W, round(W * sh / sw)
    else:
        dh, dw = H, round(H * sw / sh)
    ox, oy = (W - dw) // 2, (H - dh) // 2
    media = cfg.build / "media"; media.mkdir(exist_ok=True)
    for f in media.glob("*"):
        f.unlink()
    per, n_anim, slides = {}, 0, list(prs.slides)
    for idx, slide in enumerate(slides):
        if any(nd.tag.split('}')[-1] in ANIM for nd in slide._element.iter()):
            n_anim += 1
        vids, seen = [], set()
        for shape in slide.shapes:
            rids = {nd.get(a) for nd in shape._element.iter()
                    for a in (R + "embed", R + "link") if nd.get(a)}
            for rid in rids:
                try:
                    part = slide.part.rels[rid].target_part
                except Exception:
                    continue
                pn = str(part.partname).lower()
                ext = pn.rsplit(".", 1)[-1] if "." in pn else ""
                if ext not in VEXT or pn in seen:
                    continue
                try:
                    L, T, Wd, Hh = (int(shape.left), int(shape.top),
                                    int(shape.width), int(shape.height))
                except Exception:
                    continue
                seen.add(pn)
                fp = media / f"s{idx + 1:03d}_{rid}.{ext}"
                fp.write_bytes(part.blob)
                vids.append(Vid(fp,
                                max(0, round(ox + L / sw * dw)),
                                max(0, round(oy + T / sh * dh)),
                                max(2, round(Wd / sw * dw)),
                                max(2, round(Hh / sh * dh))))
        if vids:
            per[idx] = vids
    return per, n_anim, len(slides)


# ==================== 图片(+视频) + 音频 + 字幕 -> 段 ====================
def _ff(p) -> str:
    return str(p).replace("\\", "\\\\").replace(":", "\\:").replace("'", "\\'")


def build_segment(cfg: Config, img: Path, wav: Path, ass: Path, out: Path,
                  videos: List[Vid], total: float):
    W, H = cfg.width, cfg.height
    inp = ["-loop", "1", "-framerate", cfg.fps, "-i", img]
    for v in videos:
        inp += ["-i", v.path]
    inp += ["-i", wav]
    a_idx = 1 + len(videos)
    # 超采样下采样到输出：lanczos 抗锯齿，文字边缘锐利
    fc = [f"[0:v]scale={W}:{H}:force_original_aspect_ratio=decrease:flags=lanczos,"
          f"pad={W}:{H}:(ow-iw)/2:(oh-ih)/2:color=black,setsar=1[bg]"]
    last = "bg"
    for k, v in enumerate(videos, start=1):
        fc.append(f"[{k}:v]scale={v.w}:{v.h}:flags=lanczos,setsar=1[v{k}]")
        fc.append(f"[{last}][v{k}]overlay={v.x}:{v.y}:eof_action=repeat[t{k}]")
        last = f"t{k}"
    assf = f"ass={_ff(ass)}" + (f":fontsdir={_ff(cfg.fonts_dir)}" if cfg.fonts_dir.exists() else "")
    fc.append(f"[{last}]{assf}[vout]")
    tune = [] if videos else ["-tune", "stillimage"]   # 纯静态幻灯：x264 静态图优化更省码率/更清
    run(["ffmpeg", "-y", *inp,
         "-filter_complex", ";".join(fc),
         "-map", "[vout]", "-map", f"{a_idx}:a",
         "-af", f"loudnorm=I={cfg.loudness}:TP=-1.0:LRA=7",
         "-c:v", "libx264", "-preset", cfg.preset, *tune, "-crf", cfg.crf,
         "-pix_fmt", "yuv420p", "-r", cfg.fps,
         "-c:a", "aac", "-b:a", "192k", "-ar", "48000", "-ac", "2",
         "-t", f"{total:.3f}", out])


def concat_segments(cfg: Config, segs: List[Path], out: Path):
    listfile = cfg.build / "concat.txt"
    listfile.write_text("".join(f"file '{s.as_posix()}'\n" for s in segs))
    run(["ffmpeg", "-y", "-f", "concat", "-safe", "0", "-i", listfile, "-c", "copy", out])


# ============================== main =============================
def main():
    ap = argparse.ArgumentParser(
        description="PPTX/PDF + 中英讲稿 -> Kokoro 配音 + 双语字幕 MP4(离线·全量重建)",
        formatter_class=argparse.RawTextHelpFormatter)
    ap.add_argument("project", nargs="?", default=".", help="项目目录(含幻灯片与 talk/*.json)")
    ap.add_argument("--input", "--deck", dest="deck", default=None, help="输入 .pptx/.ppt/.pdf")
    ap.add_argument("--voice", default="am_michael", help=VOICES_HELP)
    ap.add_argument("--lang", default="a", help="a 美式 b 英式 z 中 j 日 f 法 i 意 e 西 p 葡 h 印地")
    ap.add_argument("--speed", type=float, default=1.0, help="语速(默认 1.0)")
    ap.add_argument("--effect", default="karaoke",
                    choices=["karaoke", "handwrite", "particle", "none"],
                    help="字幕特效(中英同步)：karaoke/handwrite/particle/none")
    ap.add_argument("--eng-size", type=int, default=46, dest="eng_size", help="英文字号")
    ap.add_argument("--han-size", type=int, default=32, dest="han_size", help="中文字号")
    ap.add_argument("--eng-outline", type=float, default=2.0, dest="eng_outline", help="英文粉描边粗细")
    ap.add_argument("--han-outline", type=float, default=3.0, dest="han_outline", help="中文蓝描边粗细")
    ap.add_argument("--margin-h", type=int, default=10, dest="margin_h", help="字幕左右边距")
    ap.add_argument("--eng-margin", type=int, default=4, dest="eng_margin", help="英文距底安全高度")
    ap.add_argument("--han-margin", type=int, default=4, dest="han_margin", help="中文距顶安全高度")
    ap.add_argument("--shadow", type=float, default=0.0, help="字幕阴影像素(0=无)")
    ap.add_argument("--no-shadow", action="store_true", help="(兼容旧命令)等价于 --shadow 0")
    ap.add_argument("--loudness", type=float, default=-10.0, help="统一响度 LUFS(更响可 -10)")
    ap.add_argument("--fonts-dir", default=None, dest="fonts_dir", help="字体目录(默认脚本同级 fonts/)")
    ap.add_argument("--jobs", type=int, default=8, help="并行编码线程")
    ap.add_argument("--preset", default="medium", help="x264 preset")
    ap.add_argument("--crf", type=int, default=18, help="x264 CRF(越小越清晰)")
    ap.add_argument("--dpi", type=int, default=200, help="幻灯片渲染 DPI(见文档:提DPI≈不增内存)")
    ap.add_argument("--device", default="auto", help="auto/cuda/cpu")
    a = ap.parse_args()

    cfg = Config(project=a.project, voice=a.voice, lang_code=a.lang, speed=a.speed,
                 jobs=a.jobs, preset=a.preset, crf=a.crf, dpi=a.dpi, device=a.device,
                 effect=a.effect, loudness=a.loudness,
                 eng_size=a.eng_size, han_size=a.han_size,
                 eng_outline=a.eng_outline, han_outline=a.han_outline,
                 margin_h=a.margin_h, eng_margin=a.eng_margin, han_margin=a.han_margin,
                 shadow=0.0 if a.no_shadow else a.shadow,
                 fonts_dir=Path(a.fonts_dir) if a.fonts_dir else None)

    need("ffmpeg", "sudo apt-get install -y ffmpeg")
    ensure_libass()
    cfg.build.mkdir(parents=True, exist_ok=True)

    cfg.eng_file = resolve_font(cfg.fonts_dir, cfg.eng_font, "sourcesans3", "sourcesans")
    cfg.han_file = resolve_font(cfg.fonts_dir, cfg.han_font, "sourcehansans", "hansans")
    if not (cfg.eng_file and cfg.han_file):
        sys.exit("[ERROR] fonts/ 缺少字体：需 Source Sans 3 与 Source Han Sans SC 的 OTF/TTF")
    print(f"[fonts] EN={cfg.eng_file.name} · ZH={cfg.han_file.name}")

    deck = find_deck(cfg.project, a.deck)
    images = deck_to_images(cfg, deck)
    scripts = find_scripts(cfg.project)
    if not scripts:
        sys.exit(f"[ERROR] {cfg.project / 'talk'} 下没有 .json / .txt 讲稿")
    if len(scripts) != len(images):
        sys.exit(f"[ERROR] 数量不一致：讲稿 {len(scripts)} ≠ 幻灯片 {len(images)}")

    videos_per = {}
    if deck.suffix.lower() == ".pptx":
        videos_per, n_anim, n_sl = extract_pptx_videos(cfg, deck)
        if n_anim:
            print(f"[warn] {n_anim} 页含动画/过渡，离线无法忠实还原，已忽略")
        if n_sl not in (-1, len(images)):
            print(f"[warn] PPTX 页数({n_sl})≠渲染页数({len(images)})，放弃视频叠加")
            videos_per = {}
        elif videos_per:
            tot = sum(len(v) for v in videos_per.values())
            print(f"[media] 检测到 {tot} 个内嵌视频 → 叠加到第 {sorted(k + 1 for k in videos_per)} 页")

    print(f"[layout] {cfg.width}x{cfg.height} · 中文居顶/英文居底 · 特效={cfg.effect}"
          f" · 描边 粉{cfg.eng_outline}/蓝{cfg.han_outline} · 响度{cfg.loudness}LUFS · jobs={cfg.jobs}")

    pipe = make_pipeline(cfg)
    audio_dir = cfg.build / "audio"; audio_dir.mkdir(exist_ok=True)
    sub_dir = cfg.build / "subs"; sub_dir.mkdir(exist_ok=True)
    seg_dir = cfg.build / "segments"; seg_dir.mkdir(exist_ok=True)

    segments = []
    bar = Bar(2 * len(images), "处理")

    def enc(i, img, wav, ass, seg, vids, total):
        try:
            build_segment(cfg, img, wav, ass, seg, vids, total)
        except subprocess.CalledProcessError:
            if vids:
                sys.stdout.write(f"\n[warn] 第{i}页视频叠加失败，回退静态\n")
                build_segment(cfg, img, wav, ass, seg, [], total)
            else:
                raise
        bar.update()

    with ThreadPoolExecutor(max_workers=max(1, cfg.jobs)) as ex:
        futs = []
        for i, (img, scr) in enumerate(zip(images, scripts), start=1):
            sents = load_script(scr)
            if not sents:
                raise SystemExit(f"[ERROR] 讲稿为空：{scr}")
            wav = audio_dir / f"{i:03d}.wav"
            ass = sub_dir / f"{i:03d}.ass"
            seg = seg_dir / f"{i:03d}.mp4"
            audio, cues, total = synthesize_slide(pipe, sents, cfg)       # GPU(主线程)
            write_wav(wav, audio)
            ass.write_text(build_ass(cues, cfg, total), encoding="utf-8")
            bar.update()
            futs.append(ex.submit(enc, i, img, wav, ass, seg,            # CPU(并行)
                                  videos_per.get(i - 1, []), total))
            segments.append(seg)
        for f in futs:
            f.result()

    final = cfg.project / f"{deck.stem}.mp4"
    concat_segments(cfg, segments, final)
    print(f"✅ 完成：{final}")
    shutil.rmtree(cfg.build)


if __name__ == "__main__":
    try:
        main()
    except KeyboardInterrupt:
        sys.exit("\n[中断] 用户取消")
    except subprocess.CalledProcessError as e:
        sys.exit(f"\n[ERROR] 命令失败：{' '.join(map(str, e.cmd))}\n{(e.stderr or '')[-1500:]}")